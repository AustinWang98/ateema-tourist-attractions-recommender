"""Build ChicagoDoes final capstone presentation.

Uses the original `Capstone I Initial Prosentation.pptx` as the visual template
(keeping its slide master, layouts, theme, and background/decoration imagery),
strips out every existing slide, and rebuilds the deck with the nine required
sections of content sourced from `paper/main.tex` and `paper_figures/`.

Output: ChicagoDoes_Final_Presentation.pptx (16:9, ~14 slides for a 15-min talk).
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Repo-relative paths so the script is portable after cloning.
# Layout under chicagodoes-recsys/docs/:
#   presentation/build_presentation.py      <-- this file
#   presentation/assets/                    <-- seal + formulas + screenshot
#   presentation/ChicagoDoes_Final_Presentation.pptx  <-- output
#   paper/figures/                          <-- architecture.png, mmr_diagram.png, ...
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent          # chicagodoes-recsys/
# Template lives outside the repo (Capstone I Initial Prosentation.pptx); supply via
# CHICAGODOES_TEMPLATE env var, else fall back to the previously-built deck so the
# script still produces a valid pptx from an existing copy.
import os
TEMPLATE = Path(os.environ.get(
    "CHICAGODOES_TEMPLATE",
    HERE / "ChicagoDoes_Final_Presentation.pptx",
))
OUTPUT = HERE / "ChicagoDoes_Final_Presentation.pptx"
FIGURES = ROOT / "docs" / "paper" / "figures"
ASSETS = HERE / "assets"

NAVY = RGBColor(0x27, 0x49, 0x72)
NAVY_DARK = RGBColor(0x1B, 0x2A, 0x3A)
ACCENT = RGBColor(0x4A, 0x78, 0xA5)
LIGHT_BG = RGBColor(0xE6, 0xEC, 0xF4)
TEXT_DARK = RGBColor(0x0D, 0x1B, 0x2A)
TEXT_MID = RGBColor(0x37, 0x41, 0x51)
GREY = RGBColor(0x6B, 0x72, 0x80)


# ---------------------------------------------------------------------------
# Template stripping helper
# ---------------------------------------------------------------------------
def strip_existing_slides(prs: Presentation) -> None:
    """Remove every slide already present in the template."""
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001 (private but stable)
    for sldId in list(sldIdLst):
        rId = sldId.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------
def set_text(
    tf,
    text: str | Iterable[str],
    *,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = TEXT_DARK,
    align: PP_ALIGN | None = None,
    bullet: bool = False,
) -> None:
    """Replace text frame content with given paragraphs."""
    tf.clear()
    if isinstance(text, str):
        paragraphs = [text]
    else:
        paragraphs = list(text)
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        if bullet:
            p.level = 0
        run = p.add_run()
        run.text = line
        run.font.name = "Helvetica Neue"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kwargs) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, text, **kwargs)
    return tb


def add_accent_bar(slide, left=0.0, top=0.0, width=10.0, height=0.05, color=NAVY):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_pill(slide, left, top, width, height, label, *, fill=NAVY, fg=RGBColor(0xFF, 0xFF, 0xFF), size=11):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.adjustments[0] = 0.5
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    tf = sh.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    set_text(tf, label, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return sh


def add_card(slide, left, top, width, height, title, body, *, accent=NAVY):
    """Soft-shadow card with bold title + body."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.adjustments[0] = 0.06
    box.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Accent stripe at top of card
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06)
    )
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent

    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.3), Inches(height - 0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Helvetica Neue"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Helvetica Neue"
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_MID
    return box


def add_slide_with_title(prs: Presentation, layout_idx: int, title: str | None) -> object:
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Strip every inherited placeholder text so the original template content
    # never leaks through to the new deck.
    for shape in list(slide.placeholders):
        try:
            if shape.has_text_frame:
                shape.text_frame.clear()
        except Exception:  # noqa: BLE001
            pass
    # Add our own title bar (we do NOT rely on the layout's title placeholder
    # because the template uses Google Slides-specific positioning).
    if title is not None:
        add_accent_bar(slide, 0, 0, 10, 0.04, color=NAVY)
        add_textbox(slide, 0.4, 0.18, 9.2, 0.55, title, size=22, bold=True, color=NAVY_DARK)
        add_accent_bar(slide, 0.4, 0.76, 0.6, 0.04, color=ACCENT)
    # Footer (denominator = total content slides; title + QA suppressed)
    add_textbox(slide, 0.4, 5.32, 4, 0.25, "ChicagoDoes Capstone  ·  Spring 2026",
                size=8.5, color=GREY)
    add_textbox(slide, 9.0, 5.32, 0.8, 0.25,
                f"{len(prs.slides)} / 18", size=8.5, color=GREY, align=PP_ALIGN.RIGHT)
    return slide


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Breadcrumb Only -> blank-ish
    for shape in list(slide.placeholders):
        try:
            if shape.has_text_frame:
                shape.text_frame.clear()
        except Exception:
            pass

    # Wide navy band on the left
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.4), Inches(5.62)
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY

    # University of Chicago seal (white-on-transparent) centered in the navy band
    seal_size = 1.6
    slide.shapes.add_picture(
        str(ASSETS / "uchicago_seal_white.png"),
        Inches((3.4 - seal_size) / 2), Inches(1.65),
        width=Inches(seal_size), height=Inches(seal_size),
    )

    add_textbox(slide, 0.4, 0.5, 2.6, 0.4,
                "UNIVERSITY OF CHICAGO", size=10.5, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.4, 0.85, 2.6, 0.4,
                "MS in Applied Data Science", size=10.5,
                color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    # Decorative rule under the seal
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.55), Inches(1.0), Inches(0.02)
    )
    sep.line.fill.background()
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0x9B, 0xB3, 0xCC)

    add_textbox(slide, 0.4, 3.7, 2.6, 0.3,
                "Capstone Research", size=10.5, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.4, 4.0, 2.6, 0.3,
                "Spring 2026", size=10.5,
                color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, 3.7, 1.1, 6.1, 1.2,
                "Personalized Location", size=30, bold=True, color=NAVY_DARK)
    add_textbox(slide, 3.7, 1.7, 6.1, 1.2,
                "Recommendation System", size=30, bold=True, color=NAVY_DARK)
    add_textbox(slide, 3.7, 2.3, 6.1, 0.6,
                "for ChicagoDoes", size=30, bold=True, color=ACCENT)

    add_accent_bar(slide, 3.7, 3.0, 0.9, 0.04, color=ACCENT)
    add_textbox(slide, 3.7, 3.1, 6.1, 0.4,
                "Leveraging User Interaction Data to Improve",
                size=14, italic=True, color=TEXT_MID)
    add_textbox(slide, 3.7, 3.45, 6.1, 0.4,
                "Engagement and Discovery on a Chicago Tourism Map",
                size=14, italic=True, color=TEXT_MID)

    add_textbox(slide, 3.7, 4.35, 6.1, 0.3,
                "Yiou Wang  ·  RJ Xia  ·  Kennedy Damtse",
                size=12, bold=True, color=TEXT_DARK)
    add_textbox(slide, 3.7, 4.7, 6.1, 0.3,
                "Client: Ateema / ChicagoDoes Interactive Video Maps",
                size=11, color=GREY)
    add_textbox(slide, 3.7, 5.0, 6.1, 0.3,
                "Supervisor: Don Patchell", size=11, color=GREY)
    return slide


def slide_agenda(prs):
    slide = add_slide_with_title(prs, 5, "Agenda")
    items = [
        ("1", "Introduction & Situational Analysis", 0),
        ("2", "Business Problem & Opportunity", 1),
        ("3", "Goals of Analysis", 2),
        ("4", "Deliverables", 3),
        ("5", "Literature Review", 4),
        ("6", "Data: Issues & EDA", 5),
        ("7", "Analysis Plan & Methodology", 6),
        ("8", "Model Evaluation Methodology", 7),
        ("9", "Expected Findings", 8),
    ]
    cols, rows = 3, 3
    col_w, row_h = 3.1, 1.25
    x0, y0 = 0.45, 1.15
    for num, label, i in items:
        c, r = i % cols, i // cols
        x = x0 + c * col_w
        y = y0 + r * row_h
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.55), Inches(0.55)
        )
        circle.line.fill.background()
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY
        set_text(circle.text_frame, num, size=18, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        # Label
        add_textbox(slide, x + 0.65, y + 0.05, col_w - 0.7, 0.5,
                    label, size=14, bold=True, color=NAVY_DARK)


def slide_intro(prs):
    slide = add_slide_with_title(prs, 5, "1. Introduction & Situational Analysis")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "ChicagoDoes is an interactive Chicago discovery map operated by Ateema",
                size=14, italic=True, color=ACCENT)

    cards = [
        ("The Product",
         "Interactive video map of 400+ Chicago restaurants, attractions, bars, museums, parks, theaters, and venues. Mapme-backed; reached via tourism campaigns and partner QR codes."),
        ("The Audience",
         "Tourists and locals planning city activities. Behaviorally short-session: over half of users click only 1-2 listings per session before leaving."),
        ("The Business Model",
         "Paying business partners purchase placement, CTA links, and featured marketing assets. Engagement drives partner ROI and platform retention."),
        ("The Telemetry Gap",
         "GA4 captures rich event signals (marker clicks, filters, dwell, CTA) — but the surface is fully generic. No event is operationalized into ranking."),
    ]
    col_w = 4.45
    row_h = 1.65
    for i, (t, b) in enumerate(cards):
        c, r = i % 2, i // 2
        add_card(slide, 0.4 + c * (col_w + 0.2), 1.45 + r * (row_h + 0.15),
                 col_w, row_h, t, b, accent=NAVY if c == 0 else ACCENT)


def slide_problem(prs):
    slide = add_slide_with_title(prs, 5, "2. Business Problem & Opportunity")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.45,
                "Every visitor sees the same generic map, leaving discovery and partner ROI on the table",
                size=14, italic=True, color=ACCENT)

    # Two-column problem / opportunity
    # Left: the problem
    box_l = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.55), Inches(4.55), Inches(3.6)
    )
    box_l.adjustments[0] = 0.04
    box_l.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    box_l.line.width = Pt(0.5)
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = RGBColor(0xFB, 0xE9, 0xE7)
    add_textbox(slide, 0.55, 1.65, 4.3, 0.4,
                "THE PROBLEM", size=11, bold=True,
                color=RGBColor(0xB7, 0x3A, 0x29))
    add_textbox(slide, 0.55, 2.05, 4.3, 0.4,
                "Identical experience for every visitor", size=14, bold=True, color=NAVY_DARK)
    bullets = [
        "• 15% of map-loaders click any listing marker",
        "• Most users browse only 1–2 listings per session",
        "• No personalization layer ingests GA4 events",
        "• Partner CTAs reach mismatched audiences",
        "• Each empty session loses retention + the chance to build a behavioral history",
    ]
    for i, line in enumerate(bullets):
        add_textbox(slide, 0.55, 2.5 + i * 0.45, 4.3, 0.4,
                    line, size=11, color=TEXT_MID)

    # Right: the opportunity
    box_r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.55), Inches(4.55), Inches(3.6)
    )
    box_r.adjustments[0] = 0.04
    box_r.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    box_r.line.width = Pt(0.5)
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = RGBColor(0xE6, 0xF1, 0xEC)
    add_textbox(slide, 5.25, 1.65, 4.3, 0.4,
                "THE OPPORTUNITY", size=11, bold=True,
                color=RGBColor(0x1B, 0x6B, 0x4A))
    add_textbox(slide, 5.25, 2.05, 4.3, 0.4,
                "Convert telemetry into personalization", size=14, bold=True, color=NAVY_DARK)
    bullets = [
        "• Translate GA4 events into ranking signals",
        "• Close the content-discovery gap for visitors",
        "• Increase partner CTA reach to matched users",
        "• Build a longitudinal behavioral record",
        "• Establish data infrastructure for future personalization work",
    ]
    for i, line in enumerate(bullets):
        add_textbox(slide, 5.25, 2.5 + i * 0.45, 4.3, 0.4,
                    line, size=11, color=TEXT_MID)


def slide_current_site(prs):
    slide = add_slide_with_title(prs, 5, "2. Current ChicagoDoes Experience")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "The map every visitor lands on today — a single generic surface",
                size=14, italic=True, color=ACCENT)

    # Site screenshot (1024x578, ratio ~1.77). We want it large and centered.
    pic_w = 6.3
    pic_h = pic_w * (578.0 / 1024.0)  # ~3.55
    slide.shapes.add_picture(
        str(ASSETS / "chicagodoes_site.png"),
        Inches(0.4), Inches(1.45),
        width=Inches(pic_w), height=Inches(pic_h),
    )
    # Subtle border around screenshot
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(1.45), Inches(pic_w), Inches(pic_h),
    )
    border.fill.background()
    border.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    border.line.width = Pt(0.75)
    # Caption under screenshot
    add_textbox(slide, 0.4, 1.45 + pic_h + 0.05, pic_w, 0.25,
                "Live screenshot of chicagodoes.com  ·  the recommender layer is what is missing today",
                size=9, italic=True, color=GREY)

    # Right-side observation panel
    obs_x = 0.4 + pic_w + 0.25
    obs_w = 10 - obs_x - 0.3
    add_textbox(slide, obs_x, 1.45, obs_w, 0.32,
                "What every visitor sees", size=12, bold=True, color=NAVY)
    observations = [
        "10 fixed categories in the sidebar",
        "Same default zoom & extent for all",
        "No personalized ranking — markers density-based",
        "No re-entry memory across visits",
        "Partner CTAs visible only after a marker click",
        "Search by keyword or category icon only",
    ]
    for i, o in enumerate(observations):
        y = 1.85 + i * 0.5
        # bullet dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(obs_x), Inches(y + 0.08),
            Inches(0.09), Inches(0.09),
        )
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        add_textbox(slide, obs_x + 0.2, y, obs_w - 0.2, 0.45,
                    o, size=10, color=TEXT_MID)


def slide_goals(prs):
    slide = add_slide_with_title(prs, 5, "3. Goals of Analysis")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Seven concrete objectives spanning data engineering, modeling, and validation",
                size=14, italic=True, color=ACCENT)
    goals = [
        ("Warehouse",       "Translate raw GA4 logs into a layered, queryable warehouse of user–location interactions"),
        ("Features",        "Engineer leakage-safe predictive features at user and location levels"),
        ("Hybrid Ranker",   "Combine content similarity, item/user CF, session co-visit, transitions, trending"),
        ("Cold Start",      "Handle anonymous visitors via onboarding form + K-Means behavioral archetypes"),
        ("Diversity",       "Apply MMR re-ranking so top-K lists are not category-homogeneous"),
        ("Web App",         "Ship full-stack FastAPI + JS app exposing evidence behind every recommendation"),
        ("Evaluation Plan", "Specify offline (Precision@K, NDCG, coverage) + online A/B evaluation"),
    ]
    col_w = 4.45
    row_h = 0.95
    for i, (h, body) in enumerate(goals):
        c, r = i % 2, i // 2
        x = 0.4 + c * (col_w + 0.2)
        y = 1.45 + r * row_h
        # Number badge
        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5)
        )
        num.line.fill.background()
        num.fill.solid()
        num.fill.fore_color.rgb = NAVY
        set_text(num.text_frame, str(i + 1), size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.6, y - 0.02, 3.0, 0.3,
                    h, size=12.5, bold=True, color=NAVY_DARK)
        add_textbox(slide, x + 0.6, y + 0.28, col_w - 0.6, 0.6,
                    body, size=10, color=TEXT_MID)


def slide_deliverables(prs):
    slide = add_slide_with_title(prs, 5, "4. Deliverables")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "What the client receives at the end of the engagement",
                size=14, italic=True, color=ACCENT)
    items = [
        ("BigQuery Warehouse",
         "7-layer pipeline: raw GA4 → qualified events → user–location features → candidate set. Decouples user-anonymous priors from per-user behavior to enforce leakage safety."),
        ("Hybrid Recommender",
         "Six-signal ranking engine (content + popularity + item CF + user kNN + trending + session/transition) with regime-aware weights for new vs. returning visitors."),
        ("Web Application",
         "FastAPI backend + JS frontend with evidence-backed recommendation cards (distinct users, dwell, trending) and an LLM concierge layer for itinerary planning."),
        ("Cold-Start System",
         "Onboarding form → TF-IDF pseudo-profile + K-Means archetype (Foodie Wanderer, Sightseeing Explorer, etc.). Seeds collaborative scoring for anonymous visitors."),
        ("Evaluation Methodology",
         "Documented offline plan (Precision@K, Recall@K, NDCG, coverage, diversity, MAP) + online A/B (CTR, conversion, return-visit) ready for deployment."),
        ("Code & Documentation",
         "End-to-end repository (warehouse SQL, backend Python, frontend JS) with module READMEs and a 35-page capstone paper."),
    ]
    col_w = 4.45
    row_h = 1.25
    for i, (t, b) in enumerate(items):
        c, r = i % 2, i // 2
        x = 0.4 + c * (col_w + 0.2)
        y = 1.4 + r * (row_h + 0.1)
        add_card(slide, x, y, col_w, row_h, t, b,
                 accent=NAVY if i % 2 == 0 else ACCENT)


def _lit_card_wide(slide, x, y, w, h, num, title, authors, insight, mapping, accent):
    """Wide (full-row) literature card. Header is one row: title + authors inline."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
        Inches(w), Inches(h),
    )
    card.adjustments[0] = 0.04
    card.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    card.line.width = Pt(0.5)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
        Inches(0.08), Inches(h),
    )
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent

    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.16), Inches(y + 0.10),
        Inches(0.32), Inches(0.32),
    )
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    set_text(badge.text_frame, str(num), size=11, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # title + authors (inline because the card is wide)
    add_textbox(slide, x + 0.55, y + 0.06, w * 0.55, 0.24,
                title, size=11, bold=True, color=NAVY_DARK)
    add_textbox(slide, x + 0.55 + w * 0.55, y + 0.08,
                w - 0.6 - w * 0.55, 0.22,
                authors, size=8.5, italic=True, color=GREY)

    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.55), Inches(y + 0.36),
        Inches(w - 0.7), Inches(0.012),
    )
    sep.line.fill.background()
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

    label_w = 0.95
    body_x = x + 0.18 + label_w
    body_w = w - 0.25 - label_w

    add_textbox(slide, x + 0.18, y + 0.44, label_w, 0.22,
                "Insight", size=8, bold=True, color=accent)
    add_textbox(slide, body_x, y + 0.42, body_w, 0.4,
                insight, size=9, color=TEXT_MID)

    row2_y = y + h * 0.62
    add_textbox(slide, x + 0.18, row2_y, label_w, 0.22,
                "Our system", size=8, bold=True, color=accent)
    add_textbox(slide, body_x, row2_y - 0.02, body_w, h - (row2_y - y) - 0.05,
                mapping, size=9, color=TEXT_DARK)


def _lit_card_compact(slide, x, y, w, h, num, title, authors, insight, mapping, accent):
    """Compact card for a 2x2 grid. Title and authors stack vertically."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
        Inches(w), Inches(h),
    )
    card.adjustments[0] = 0.04
    card.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    card.line.width = Pt(0.5)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
        Inches(0.08), Inches(h),
    )
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent

    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.16), Inches(y + 0.10),
        Inches(0.32), Inches(0.32),
    )
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    set_text(badge.text_frame, str(num), size=11, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    add_textbox(slide, x + 0.55, y + 0.06, w - 0.7, 0.26,
                title, size=11, bold=True, color=NAVY_DARK)
    add_textbox(slide, x + 0.55, y + 0.32, w - 0.7, 0.22,
                authors, size=8.5, italic=True, color=GREY)

    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.55), Inches(y + 0.58),
        Inches(w - 0.7), Inches(0.012),
    )
    sep.line.fill.background()
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

    # Non-overlapping insight + mapping regions (8.5pt body for compact fit)
    insight_label_y = y + 0.60
    insight_body_y = y + 0.79
    insight_body_h = 0.42
    mapping_label_y = insight_body_y + insight_body_h + 0.02
    mapping_body_y = mapping_label_y + 0.18
    mapping_body_h = max(0.30, y + h - 0.05 - mapping_body_y)

    add_textbox(slide, x + 0.18, insight_label_y, w - 0.36, 0.18,
                "Insight", size=8, bold=True, color=accent)
    add_textbox(slide, x + 0.18, insight_body_y, w - 0.36, insight_body_h,
                insight, size=8.5, color=TEXT_MID)

    add_textbox(slide, x + 0.18, mapping_label_y, w - 0.36, 0.18,
                "Our system", size=8, bold=True, color=accent)
    add_textbox(slide, x + 0.18, mapping_body_y, w - 0.36, mapping_body_h,
                mapping, size=8.5, color=TEXT_DARK)


def slide_litreview_1(prs):
    slide = add_slide_with_title(prs, 5, "5. Literature Review — Foundations (1/2)")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Core recommender-systems theory underpinning our scoring blend",
                size=13, italic=True, color=ACCENT)
    entries = [
        (1,
         "Implicit-Feedback Collaborative Filtering",
         "Hu, Koren & Volinsky (2008)  ·  Koren, Bell & Volinsky (2009)",
         "Treat clicks / dwell as confidence-weighted positive signals; matrix factorization with confidence is the canonical implicit-feedback model.",
         "Foundation for our item–item and user–user CF terms. We use confidence-weighted co-occurrence (sparse Jaccard) rather than full MF, because the GA4 export is short and the catalog is small (~400 locations)."),
        (2,
         "Content-Based Methods & TF-IDF",
         "Salton & Buckley (1988)  ·  Pazzani & Billsus (2007)",
         "Item attributes vectorized via TF-IDF + cosine similarity outperform raw category matching on sparse catalogs.",
         "We TF-IDF-vectorize each location's category tokens, then compute cosine similarity to a TF-IDF user profile (history-weighted). This is also the backbone of the cold-start pseudo-profile."),
        (3,
         "Hybrid Recommenders",
         "Burke (2002)  ·  Çano & Morisio (2017)",
         "A weighted/switching hybrid of content + CF empirically dominates either paradigm alone on sparse, implicit-feedback data — the regime where ChicagoDoes operates.",
         "Directly motivates our six-signal weighted blend, with different weight vectors for new-visitor vs. returning-user regimes to handle the cold-start↔warm-start transition."),
    ]
    card_h = 1.18
    for i, e in enumerate(entries):
        y = 1.40 + i * (card_h + 0.08)
        _lit_card_wide(slide, 0.4, y, 9.25, card_h,
                       e[0], e[1], e[2], e[3], e[4],
                       accent=NAVY if i % 2 == 0 else ACCENT)


def slide_litreview_2(prs):
    slide = add_slide_with_title(prs, 5, "5. Literature Review — Applied & Domain (2/2)")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Session, diversification, tourism, and LLM literature shaping the deployment surface",
                size=13, italic=True, color=ACCENT)
    entries = [
        (4,
         "Session-Aware Recommendation",
         "Hidasi et al. (2016)  ·  Quadrana et al. (2018)",
         "Anonymous, short-session traffic still carries in-session intent — captured via co-visit graphs and previous→next transition counts.",
         "Same-session Jaccard (0.6) + transition probability (0.4). Crucial since >50% of ChicagoDoes traffic is single-session."),
        (5,
         "Diversity — Maximal Marginal Relevance",
         "Carbonell & Goldstein (1998)",
         "Relevance-only top-K lists become category-homogeneous; MMR re-ranks by trading λ relevance vs. (1−λ) novelty.",
         "Applied at λ = 0.7 over the top-60 candidates — prevents the 'eight bars in a row' failure mode in baseline runs."),
        (6,
         "Tourism RS & Cold Start",
         "Borrás et al. (2014)  ·  Gavalas et al. (2014)",
         "Tourism is dominated by cold-start anonymous visitors; standard recipe = lightweight onboarding form + behavioral segmentation.",
         "Onboarding form → TF-IDF pseudo-profile + K-Means archetype (k = 6: Foodie, Sightseer, …) seeding CF for new users."),
        (7,
         "LLMs as Presentation Layer",
         "Lewis et al. (2020)  ·  Wu et al. (2024)",
         "LLMs perform best when ranking is deterministic and the LLM only narrates / re-orders within a small retrieved pool (RAG pattern).",
         "ChicagoDoes 'concierge': deterministic hybrid produces the pool; LLM writes copy + ordering — never invents locations."),
    ]
    # 2x2 compact grid: each card 4.55" wide × 1.85" tall
    card_w = 4.55
    card_h = 1.85
    for i, e in enumerate(entries):
        c, r = i % 2, i // 2
        x = 0.4 + c * (card_w + 0.15)
        y = 1.40 + r * (card_h + 0.10)
        _lit_card_compact(slide, x, y, card_w, card_h,
                          e[0], e[1], e[2], e[3], e[4],
                          accent=NAVY if i % 2 == 0 else ACCENT)


def slide_data(prs):
    slide = add_slide_with_title(prs, 5, "6. Data: Sources & Issues")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "GA4 event stream + canonical dimensions, layered into 7 warehouse tables",
                size=14, italic=True, color=ACCENT)

    # Left: Sources
    add_textbox(slide, 0.4, 1.5, 4.5, 0.35,
                "DATA SOURCES", size=12, bold=True, color=NAVY)
    sources = [
        ("GA4 Daily Events",  "events_YYYYMMDD + events_intraday tables in BigQuery"),
        ("Location Dimension", "Official Mapme / Ateema location & category lists"),
        ("Crawler Bridge",    "Cleaned location × category mapping"),
        ("Engagement Filter", "Strong / weak / nav / noise event classification"),
    ]
    for i, (h, b) in enumerate(sources):
        y = 1.85 + i * 0.78
        add_pill(slide, 0.4, y, 1.5, 0.32, h, fill=NAVY, size=10)
        add_textbox(slide, 2.0, y, 2.9, 0.35, b, size=10.5, color=TEXT_MID)

    # Right: Data Issues
    add_textbox(slide, 5.1, 1.5, 4.5, 0.35,
                "KNOWN DATA ISSUES", size=12, bold=True,
                color=RGBColor(0xB7, 0x3A, 0x29))
    issues = [
        ("Implicit only",        "No star ratings or post-visit feedback exist"),
        ("Short window",         "Only a few weeks of GA4 export — no annual seasonality"),
        ("Anonymous users",      "Most visitors single-session; no demographics"),
        ("Conversion blind spot","Once user clicks off the map, downstream events are not observable"),
        ("Sparse content",       "Locations carry categories but few free-text descriptions"),
        ("Leakage risk",         "Aggregated user features could leak labels → isolated"),
    ]
    for i, (h, b) in enumerate(issues):
        y = 1.85 + i * 0.5
        add_pill(slide, 5.1, y, 1.45, 0.32, h,
                 fill=RGBColor(0xB7, 0x3A, 0x29), size=9.5)
        add_textbox(slide, 6.65, y, 2.95, 0.32, b, size=9.5, color=TEXT_MID)


def slide_eda_engagement(prs):
    slide = add_slide_with_title(prs, 5, "6. EDA — Engagement Is Shallow & Heavy-Tailed")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.35,
                "Funnel drops sharply at first marker click; most users click only 1–2 listings",
                size=12, italic=True, color=ACCENT)

    # Left: funnel chart, sized to fit
    slide.shapes.add_picture(
        str(FIGURES / "fig_engagement_funnel.png"),
        Inches(0.3), Inches(1.4),
        width=Inches(4.7),
    )
    # Right: behavior depth (already a 2-panel chart)
    slide.shapes.add_picture(
        str(FIGURES / "fig_behavior_depth.png"),
        Inches(5.1), Inches(1.7),
        width=Inches(4.8),
    )

    # Bottom takeaway band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(4.75),
        Inches(9.25), Inches(0.5),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    add_textbox(slide, 0.55, 4.78, 1.7, 0.45,
                "Key takeaway", size=11, bold=True, color=NAVY)
    add_textbox(slide, 2.2, 4.78, 7.4, 0.45,
                "Only 15% of map-loaders click a marker and >50% click ≤2 listings per session — the recommender must front-load the most relevant places in the first two impressions.",
                size=10, color=TEXT_MID)


def slide_eda_listings(prs):
    slide = add_slide_with_title(prs, 5, "6. EDA — Repeat-Driven Popularity & Multi-Intent Visitors")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.35,
                "Same listings dominate by total clicks vs. unique users; first-click entry points are diverse",
                size=12, italic=True, color=ACCENT)

    # Left: side-by-side listings ranking (clicks vs unique users)
    slide.shapes.add_picture(
        str(FIGURES / "fig_listing_engagement.png"),
        Inches(0.3), Inches(1.4),
        width=Inches(5.4),
    )
    # Right: entry-point chart
    slide.shapes.add_picture(
        str(FIGURES / "fig_entry_points.png"),
        Inches(5.8), Inches(1.4),
        width=Inches(4.0),
    )

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(4.75),
        Inches(9.25), Inches(0.5),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    add_textbox(slide, 0.55, 4.78, 2.0, 0.45,
                "Design implications", size=11, bold=True, color=NAVY)
    add_textbox(slide, 2.6, 4.78, 7.0, 0.45,
                "Popularity must be computed per distinct user (not raw clicks) — and entry-point heterogeneity justifies segmenting visitors into K-Means behavioral archetypes.",
                size=10, color=TEXT_MID)


def slide_methodology(prs):
    slide = add_slide_with_title(prs, 5, "7. Analysis Plan — System Architecture")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Layered warehouse → six-signal hybrid ranker → MMR diversification → web app",
                size=13, italic=True, color=ACCENT)
    # architecture.png is 1879x1099 (ratio 1.71). To fit ≤3.5" tall use width 6.0"
    pic_h = 3.55
    pic_w = pic_h * (1879.0 / 1099.0)
    slide.shapes.add_picture(
        str(FIGURES / "architecture.png"),
        Inches((10 - pic_w) / 2), Inches(1.4),
        width=Inches(pic_w), height=Inches(pic_h),
    )
    # Bottom note
    add_textbox(slide, 0.4, 5.05, 9.2, 0.3,
                "Each ranking signal is leakage-safe (user-anonymous priors decoupled from per-user behavior)",
                size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def slide_score_blend(prs):
    slide = add_slide_with_title(prs, 5, "7. Modeling Framework — Six-Signal Score Blend")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Each location scored on six normalized signals; weights switch by user regime",
                size=13, italic=True, color=ACCENT)

    # Left: 6 signals as compact cards
    signals = [
        ("Content sim.",    "TF-IDF cosine between user profile and location categories"),
        ("Popularity",      "Distinct-user-based prior; leakage-safe"),
        ("Item–item CF",    "Sparse Jaccard between locations engaged by overlapping users"),
        ("User–user kNN",   "Top-5 nearest users by category-share cosine"),
        ("Trending",        "Recent vs. early-window event ratio per location"),
        ("Session/transit", "Same-session Jaccard (0.6) + previous→next transition (0.4)"),
    ]
    add_textbox(slide, 0.4, 1.45, 5.0, 0.32,
                "Six scoring signals", size=11, bold=True, color=NAVY)
    for i, (h, b) in enumerate(signals):
        y = 1.78 + i * 0.55
        add_pill(slide, 0.4, y, 1.55, 0.32, h, fill=NAVY, size=9.5)
        add_textbox(slide, 2.05, y, 3.25, 0.4, b, size=9.5, color=TEXT_MID)

    # Right: regime weights chart from paper
    add_textbox(slide, 5.5, 1.45, 4.2, 0.32,
                "Weights by user regime", size=11, bold=True, color=NAVY)
    slide.shapes.add_picture(
        str(FIGURES / "score_blend.png"),
        Inches(5.5), Inches(1.78),
        width=Inches(4.2),
    )


def _add_formula_png(slide, png_path, *, center_x, center_y, max_w, max_h):
    """Place a formula PNG, scaled to fit within (max_w, max_h) and centered on
    (center_x, center_y). All units in inches."""
    from PIL import Image
    with Image.open(png_path) as im:
        ratio = im.size[0] / im.size[1]
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    slide.shapes.add_picture(
        str(png_path),
        Inches(center_x - w / 2), Inches(center_y - h / 2),
        width=Inches(w), height=Inches(h),
    )


def slide_algorithm_detail(prs):
    slide = add_slide_with_title(prs, 5, "7. Algorithm Detail — Scoring & Re-Ranking")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Composite location score, signal computations, and MMR diversification",
                size=13, italic=True, color=ACCENT)

    # ---- Top: master composite score formula band ----
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.45), Inches(9.25), Inches(0.95),
    )
    band.adjustments[0] = 0.08
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    add_textbox(slide, 0.55, 1.55, 9.0, 0.3,
                "Composite location score for user u and location ℓ",
                size=10.5, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1))
    # Composite formula PNG centered in band
    _add_formula_png(
        slide,
        ASSETS / "formulas" / "composite_white.png",
        center_x=5.0, center_y=2.10, max_w=8.8, max_h=0.45,
    )

    # ---- Three formula cards (left, middle, right) ----
    card_meta = [
        ("Content similarity (TF-IDF)",
         "content_sim",
         "v_u, v_ℓ are L2-normalized TF-IDF vectors over category tokens. v_u is the user's history-weighted average."),
        ("Item–item CF (Jaccard)",
         "item_cf",
         "U_i = set of users who engaged location i. Per-user score = max Jaccard against history. Sparse and leakage-safe."),
        ("Trending (recency lift)",
         "trending",
         "Recent vs. early strong-event rates per location. Min-max normalized to [0,1] across the catalog."),
    ]
    card_w = 2.95
    card_h = 1.65
    card_y = 2.55
    for i, (h, formula_key, body) in enumerate(card_meta):
        x = 0.4 + i * (card_w + 0.15)
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.06
        card.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Accent strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(card_y),
            Inches(card_w), Inches(0.06),
        )
        strip.line.fill.background()
        strip.fill.solid()
        strip.fill.fore_color.rgb = NAVY if i % 2 == 0 else ACCENT
        # Title
        add_textbox(slide, x + 0.12, card_y + 0.1, card_w - 0.24, 0.3,
                    h, size=11.5, bold=True, color=NAVY_DARK)
        # Formula PNG inside a light-blue box
        fb_y = card_y + 0.45
        fb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.12), Inches(fb_y),
            Inches(card_w - 0.24), Inches(0.55),
        )
        fb.line.fill.background()
        fb.fill.solid()
        fb.fill.fore_color.rgb = LIGHT_BG
        _add_formula_png(
            slide,
            ASSETS / "formulas" / f"{formula_key}.png",
            center_x=x + card_w / 2,
            center_y=fb_y + 0.275,
            max_w=card_w - 0.36,
            max_h=0.48,
        )
        # Body
        add_textbox(slide, x + 0.12, card_y + 1.05, card_w - 0.24, 0.55,
                    body, size=9.5, color=TEXT_MID)

    # ---- MMR formula band at the bottom ----
    mmr_band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(4.3), Inches(9.25), Inches(0.95),
    )
    mmr_band.adjustments[0] = 0.10
    mmr_band.line.color.rgb = ACCENT
    mmr_band.line.width = Pt(0.75)
    mmr_band.fill.solid()
    mmr_band.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFB)
    add_textbox(slide, 0.55, 4.35, 9.0, 0.28,
                "MMR diversification (λ = 0.7) — applied to top-60 candidates",
                size=10.5, bold=True, color=NAVY)
    _add_formula_png(
        slide,
        ASSETS / "formulas" / "mmr.png",
        center_x=5.0, center_y=4.92, max_w=8.6, max_h=0.55,
    )
    add_textbox(slide, 0.55, 5.20, 9.0, 0.18,
                "S = items already in the result list  ·  70% relevance, 30% novelty vs. already-picked items",
                size=8.5, italic=True, color=TEXT_MID, align=PP_ALIGN.CENTER)


def slide_coldstart_mmr(prs):
    slide = add_slide_with_title(prs, 5, "7. Cold Start + Diversity Re-Ranking")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Form-derived pseudo-profile + K-Means archetypes, then MMR diversifies the top-K",
                size=13, italic=True, color=ACCENT)

    add_textbox(slide, 0.4, 1.4, 4.5, 0.32,
                "Cold-start handling", size=12, bold=True, color=NAVY)
    cs_steps = [
        ("Onboarding form",
         "Interests, traveler type, vibe → mapped to category weights via lookup tables"),
        ("TF-IDF projection",
         "Form weights projected into the same TF-IDF space as locations → pseudo-profile"),
        ("K-Means archetype",
         "k=6 segments: Foodie Wanderer, Sightseeing Explorer, Cultural Enthusiast, etc."),
        ("Seed set",
         "Archetype-popular locations seed the item–item CF score for cold-start users"),
    ]
    for i, (h, b) in enumerate(cs_steps):
        y = 1.78 + i * 0.78
        add_pill(slide, 0.4, y, 1.55, 0.32, h, fill=NAVY, size=9.5)
        add_textbox(slide, 2.05, y, 2.95, 0.55, b, size=9.5, color=TEXT_MID)

    # Right: MMR diagram
    add_textbox(slide, 5.4, 1.4, 4.4, 0.32,
                "MMR diversification (λ = 0.7)", size=12, bold=True, color=NAVY)
    slide.shapes.add_picture(
        str(FIGURES / "mmr_diagram.png"),
        Inches(5.4), Inches(1.78),
        width=Inches(4.4),
    )
    add_textbox(slide, 5.4, 4.6, 4.4, 0.5,
                "70% relevance + 30% distance from already-picked items → re-shuffles homogeneous top-K lists without dropping the relevance leader",
                size=9.5, italic=True, color=TEXT_MID)


def slide_eval(prs):
    slide = add_slide_with_title(prs, 5, "8. Model Evaluation Methodology")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Offline + online + assumption checks — scoped for the next deployment phase",
                size=13, italic=True, color=ACCENT)

    cards = [
        ("Offline Holdout",
         "Leave-last-N per user; time-stratified split since trending and session signals are time-aware. Train on remaining interactions."),
        ("Ranking Metrics",
         "Precision@K, Recall@K, Hit Rate@K, NDCG@K, MAP at K ∈ {5, 10, 20}"),
        ("List Quality",
         "Catalog coverage, intra-list cosine diversity (w/ vs w/o MMR), novelty, category entropy"),
        ("Baselines",
         "(i) Global popularity  (ii) Random within interest  (iii) Content-only TF-IDF  (iv) Item–item CF only"),
        ("Ablations",
         "Remove each scoring term in turn → quantify marginal contribution; re-tune regime weights"),
        ("Online A/B",
         "Behind feature flag: locations/session, CTR, dwell, CTA conversions, return-visit rate, partner-side impressions"),
    ]
    col_w = 4.45
    row_h = 1.15
    for i, (t, b) in enumerate(cards):
        c, r = i % 2, i // 2
        x = 0.4 + c * (col_w + 0.2)
        y = 1.45 + r * (row_h + 0.12)
        add_card(slide, x, y, col_w, row_h, t, b,
                 accent=NAVY if i % 2 == 0 else ACCENT)


def slide_findings(prs):
    slide = add_slide_with_title(prs, 5, "9. Expected Findings")
    add_textbox(slide, 0.4, 0.95, 9.2, 0.4,
                "Mapped to the planned baselines + observed in design-validation runs",
                size=13, italic=True, color=ACCENT)

    rows = [
        ("Precision@K, NDCG@K",  "Hybrid > content-only > collab-only > popularity",
         "Content backstop + CF refinement in same model"),
        ("Hit Rate@K (cold-start)", "Hybrid ≈ popularity",
         "Popularity is a strong baseline for new visitors"),
        ("Intra-list diversity", "Hybrid + MMR ≫ Hybrid w/o MMR",
         "MMR explicitly trades 30% relevance for diversity"),
        ("Catalog coverage",     "Hybrid ≫ popularity",
         "CF + content surface the long tail"),
        ("Engagement filter",    "Top-K Jaccard ≥ 0.7 across policies",
         "Filter removes noise without destroying signal"),
        ("Online CTR",           "Personalization wins, partner CTA clicks lift",
         "Standard uplift pattern in DMO / tourism literature"),
    ]
    # Header
    header_y = 1.45
    add_pill(slide, 0.4, header_y, 2.6, 0.34, "Dimension", fill=NAVY, size=10)
    add_pill(slide, 3.1, header_y, 3.2, 0.34, "Expectation", fill=NAVY, size=10)
    add_pill(slide, 6.4, header_y, 3.25, 0.34, "Reasoning", fill=NAVY, size=10)

    for i, (a, b, c) in enumerate(rows):
        y = header_y + 0.42 + i * 0.52
        # Zebra striping
        if i % 2 == 0:
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y - 0.04),
                Inches(9.25), Inches(0.5),
            )
            stripe.line.fill.background()
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = LIGHT_BG
        add_textbox(slide, 0.5, y, 2.55, 0.45, a, size=10, bold=True, color=NAVY_DARK)
        add_textbox(slide, 3.2, y, 3.1, 0.45, b, size=10, color=TEXT_DARK)
        add_textbox(slide, 6.5, y, 3.2, 0.45, c, size=10, color=TEXT_MID)


def slide_qa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in list(slide.placeholders):
        try:
            if shape.has_text_frame:
                shape.text_frame.clear()
        except Exception:
            pass

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.62)
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY

    add_textbox(slide, 0.5, 1.8, 9.0, 0.8,
                "Thank You", size=44, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 2.7, 9.0, 0.5,
                "Questions & Discussion", size=22,
                color=RGBColor(0xCB, 0xD5, 0xE1),
                align=PP_ALIGN.CENTER)

    add_accent_bar(slide, 4.0, 3.35, 2.0, 0.04, color=ACCENT)

    add_textbox(slide, 0.5, 3.6, 9.0, 0.4,
                "Yiou Wang  ·  RJ Xia  ·  Kennedy Damtse",
                size=14, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.0, 9.0, 0.4,
                "Supervisor: Don Patchell",
                size=12, color=RGBColor(0xCB, 0xD5, 0xE1),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.35, 9.0, 0.4,
                "Client: Ateema / ChicagoDoes Interactive Video Maps",
                size=12, color=RGBColor(0xCB, 0xD5, 0xE1),
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main() -> None:
    prs = Presentation(TEMPLATE)
    strip_existing_slides(prs)

    slide_title(prs)
    slide_agenda(prs)
    slide_intro(prs)
    slide_problem(prs)
    slide_current_site(prs)
    slide_goals(prs)
    slide_deliverables(prs)
    slide_litreview_1(prs)
    slide_litreview_2(prs)
    slide_data(prs)
    slide_eda_engagement(prs)
    slide_eda_listings(prs)
    slide_methodology(prs)
    slide_score_blend(prs)
    slide_algorithm_detail(prs)
    slide_coldstart_mmr(prs)
    slide_eval(prs)
    slide_findings(prs)
    slide_qa(prs)

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
